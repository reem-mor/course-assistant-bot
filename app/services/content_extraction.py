"""Text extraction from Drive materials.

Course materials are mostly binary (pptx/pdf/docx), so "Drive native extraction" only
covers Google-native docs; binary files are downloaded and parsed locally. Dispatch is by
mime type. A size cap avoids pulling very large files into memory.
"""

from __future__ import annotations

import io

from app.core.logging import get_logger
from app.services.drive import DriveService
from app.services.drive_parser import file_id, file_mime, file_name, file_size

_log = get_logger("content_extraction")

# Skip extraction for files larger than this (links are offered instead).
_MAX_EXTRACT_BYTES = 20 * 1024 * 1024

_GOOGLE_NATIVE_PREFIX = "application/vnd.google-apps"
_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _matches_ext(name: str, *exts: str) -> bool:
    return name.lower().endswith(exts)


async def extract_text(drive: DriveService, file: dict[str, object]) -> str:
    """Extract plain text from a Drive file dict, dispatching by mime/extension.

    Returns an empty string for unsupported or oversized files (the caller falls back to
    links rather than content).
    """
    mime = file_mime(file)
    name = file_name(file)
    fid = file_id(file)

    if mime.startswith(_GOOGLE_NATIVE_PREFIX):
        return await drive.read_file_content(fid)

    size = file_size(file)
    if size is not None and size > _MAX_EXTRACT_BYTES:
        _log.info("extract_skipped_large", name=name, size=size)
        return ""

    try:
        data = await drive.download_file(fid)
        if mime == _PDF_MIME or _matches_ext(name, ".pdf"):
            return _extract_pdf(data)
        if mime == _DOCX_MIME or _matches_ext(name, ".docx"):
            return _extract_docx(data)
        if mime == _PPTX_MIME or _matches_ext(name, ".pptx"):
            return _extract_pptx(data)
    except Exception:  # extraction must never crash a summary request
        _log.exception("extract_failed", name=name, mime=mime)
        return ""
    return ""
